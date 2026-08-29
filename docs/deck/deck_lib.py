"""
deck_lib.py — Watheeq investor-deck design system (python-pptx).

A house component library so each slide is composed from consistent, premium
primitives + diagram archetypes (flows, arrows, stacks, hub-spoke, tables,
quadrants, timelines, KPI tiles, capability maps, gauges, screenshot frames).

Slide size: 13.333 x 7.5 in (16:9).  All coords are in INCHES (floats).
Colors are HEX strings WITHOUT '#'.  Use rgb('C6A962') -> RGBColor.

Author note for contributors: compose the COMPONENTS below; do not hand-roll raw
python-pptx unless a component is missing. Every public component returns a
bbox dict {x,y,w,h} of the area it occupied so callers can stack/align.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn, nsdecls
from pptx.oxml import parse_xml
import re, os, math, struct, subprocess, glob

EMU_IN = 914400
SW, SH = 13.333, 7.5

# ----------------------------------------------------------------- palette
class C:
    ink      = '15211F'   # near-black text
    ink_soft = '2A3A38'
    teal_dk  = '0A3D3D'   # cover / primary dark
    teal_dk2 = '063031'   # darker
    teal     = '0D4B4F'   # brand teal (DataStream)
    teal_md  = '14686E'
    green    = '1B5E20'   # leaf (delivered accent)
    gold     = 'C6A962'   # brand gold
    gold_dk  = 'A68B42'
    gold_lt  = 'E3D2A8'
    gold_pale= 'F3EAD2'
    navy     = '0033A1'   # brand navy (platform tie / data viz)
    navy_lt  = '1D5BFF'
    muted    = '5B6B6B'
    muted_lt = '8A9896'
    line     = 'D7E2DE'
    line_dk  = 'C2D2CD'
    card     = 'F2F7F5'
    card2    = 'E9F1EE'
    card3    = 'DCEAE5'
    bg       = 'FBFCFC'
    white    = 'FFFFFF'
    success  = '1F9D6B'
    warn     = 'D79A2B'
    crit     = 'C24233'
    info     = '2E86C1'

HEAD = 'Georgia'      # editorial serif for headlines (ubiquitous Win/Mac)
BODY = 'Calibri'      # clean sans for body
NUMS = 'Georgia'      # big numerals

def rgb(h):
    return RGBColor.from_string(h.lstrip('#'))

# asset paths
_ASSETS = os.path.join(os.path.dirname(os.path.abspath(__file__)), '..', '..')
LOGO_ICON = '/private/tmp/claude-501/-Users-mac-clario360/533e7b8c-6306-4996-9a65-34aeea71f04d/scratchpad/assets/icon.png'
LOGO_WORD = '/private/tmp/claude-501/-Users-mac-clario360/533e7b8c-6306-4996-9a65-34aeea71f04d/scratchpad/assets/logo.png'
SHOTS_DIR = '/Users/mac/clario360/.artifacts/lex-browser-evidence-2026-06-25T07-42-33-357Z/screenshots'

def shot(name):
    p = os.path.join(SHOTS_DIR, name if name.endswith('.png') else name + '.png')
    return p if os.path.exists(p) else None

# ----------------------------------------------------------------- low-level
def _spPr(shape):
    return shape._element.spPr

def no_shadow(shape):
    shape.shadow.inherit = False
    return shape

def soft_shadow(shape, blur=0.11, dist=0.05, direction=5400000, hexv='10201E', alpha=34):
    sp = _spPr(shape)
    for e in sp.findall(qn('a:effectLst')):
        sp.remove(e)
    xml = ('<a:effectLst %s><a:outerShdw blurRad="%d" dist="%d" dir="%d" rotWithShape="0">'
           '<a:srgbClr val="%s"><a:alpha val="%d"/></a:srgbClr></a:outerShdw></a:effectLst>'
           ) % (nsdecls('a'), int(blur*EMU_IN), int(dist*EMU_IN), direction, hexv, int(alpha*1000))
    sp.append(parse_xml(xml))
    return shape

def fill_alpha(shape, hexv, alpha_pct):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(hexv)
    sf = _spPr(shape).find(qn('a:solidFill'))
    srgb = sf.find(qn('a:srgbClr'))
    srgb.append(parse_xml('<a:alpha %s val="%d"/>' % (nsdecls('a'), int(alpha_pct*1000))))
    return shape

def gradient_fill(shape, hex1, hex2, angle_deg=90):
    try:
        shape.fill.gradient()
        gs = shape.fill.gradient_stops
        gs[0].color.rgb = rgb(hex1)
        gs[-1].color.rgb = rgb(hex2)
        try:
            shape.fill.gradient_angle = angle_deg
        except Exception:
            pass
    except Exception:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(hex1)
    return shape

def rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0, radius=None,
         shape=MSO_SHAPE.RECTANGLE, shadow=False, dash=None):
    if radius is not None:
        shape = MSO_SHAPE.ROUNDED_RECTANGLE
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    no_shadow(sp)
    if radius is not None:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = rgb(fill)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = rgb(line); sp.line.width = Pt(line_w)
        if dash:
            ln = sp.line._get_or_add_ln()
            d = parse_xml('<a:prstDash %s val="%s"/>' % (nsdecls('a'), dash))
            ln.append(d)
    if shadow:
        soft_shadow(sp)
    return sp

def oval(slide, x, y, w, h, fill=None, line=None, line_w=1.0, shadow=False):
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    no_shadow(sp)
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = rgb(fill)
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = rgb(line); sp.line.width = Pt(line_w)
    if shadow: soft_shadow(sp)
    return sp

def freeshape(slide, mso, x, y, w, h, fill=None, line=None, line_w=1.0):
    sp = slide.shapes.add_shape(mso, Inches(x), Inches(y), Inches(w), Inches(h))
    no_shadow(sp)
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = rgb(fill)
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = rgb(line); sp.line.width = Pt(line_w)
    return sp

def poly(slide, pts_in, fill=None, line=None, line_w=1.0):
    """Freeform polygon from list of (x,y) inch points. Text-safe (upright)."""
    emu = [(int(px*EMU_IN), int(py*EMU_IN)) for px, py in pts_in]
    fb = slide.shapes.build_freeform(emu[0][0], emu[0][1], scale=1.0)
    fb.add_line_segments(emu[1:], close=True)
    sp = fb.convert_to_shape()
    no_shadow(sp)
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = rgb(fill)
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = rgb(line); sp.line.width = Pt(line_w)
    return sp

def _arrowhead(line, end='tail', kind='triangle', size='med'):
    ln = line._get_or_add_ln()
    tag = 'a:tailEnd' if end == 'tail' else 'a:headEnd'
    for e in ln.findall(qn(tag)):
        ln.remove(e)
    el = parse_xml('<%s %s type="%s" w="%s" len="%s"/>' % (tag, nsdecls('a'), kind, size, size))
    ln.append(el)

def connector(slide, x1, y1, x2, y2, color='5B6B6B', w=1.6, kind='straight',
              head=True, tail=False, dash=None):
    mso = MSO_CONNECTOR.ELBOW if kind == 'elbow' else MSO_CONNECTOR.STRAIGHT
    cn = slide.shapes.add_connector(mso, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = rgb(color)
    cn.line.width = Pt(w)
    no_shadow(cn)
    if dash:
        ln = cn.line._get_or_add_ln()
        ln.append(parse_xml('<a:prstDash %s val="%s"/>' % (nsdecls('a'), dash)))
    if head: _arrowhead(cn.line, 'tail', 'triangle', 'med')
    if tail: _arrowhead(cn.line, 'head', 'triangle', 'med')
    return cn

# ----------------------------------------------------------------- text
def R(text, size, color, bold=False, italic=False, font=BODY, spacing=None):
    return dict(t=str(text), s=size, c=color, b=bold, i=italic, f=font, sp=spacing)

def _apply_runs(p, runs):
    for rspec in runs:
        run = p.add_run()
        run.text = rspec['t']
        run.font.size = Pt(rspec['s'])
        run.font.color.rgb = rgb(rspec['c'])
        run.font.name = rspec['f']
        run.font.bold = rspec['b']
        run.font.italic = rspec['i']
        if rspec.get('sp') is not None:
            # letter spacing in points -> EMU? use spc attribute (1/100 pt)
            rPr = run._r.get_or_add_rPr()
            rPr.set('spc', str(int(rspec['sp'] * 100)))

def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    return tf, tb

def add_para(tf, runs, align=PP_ALIGN.LEFT, space_after=0, space_before=0,
             line_spacing=1.0, first=False):
    if first and not tf.paragraphs[0].runs:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = align
    if space_after: p.space_after = Pt(space_after)
    if space_before: p.space_before = Pt(space_before)
    if line_spacing: p.line_spacing = line_spacing
    _apply_runs(p, runs)
    return p

def label(slide, x, y, w, h, text, size, color, bold=False, italic=False,
          font=BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=None,
          line_spacing=1.0):
    tf, tb = textbox(slide, x, y, w, h, anchor=anchor)
    add_para(tf, [R(text, size, color, bold, italic, font, spacing)],
             align=align, line_spacing=line_spacing, first=True)
    return tf

def bullets(slide, x, y, w, h, items, size=16, color=None, marker=C.gold,
            gap=10, line_spacing=1.06, marker_char='▪'):
    color = color or C.ink
    tf, tb = textbox(slide, x, y, w, h)
    for i, it in enumerate(items):
        if isinstance(it, tuple):
            head, sub = it
        else:
            head, sub = it, None
        runs = [R(marker_char + '  ', size, marker, bold=True), R(head, size, color)]
        add_para(tf, runs, space_after=(2 if sub else gap), line_spacing=line_spacing, first=(i == 0))
        if sub:
            add_para(tf, [R('     ' + sub, size - 3, C.muted)], space_after=gap, line_spacing=1.0)
    return tf

# ----------------------------------------------------------------- chips / badges
def chip(slide, x, y, w, h, text, fill=C.teal_dk, txt=C.white, size=12, bold=True,
         radius=0.5, font=BODY, line=None):
    sp = rect(slide, x, y, w, h, fill=fill, radius=radius, line=line, line_w=1.0)
    tf = sp.text_frame; tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
    tf.margin_top = 0; tf.margin_bottom = 0
    add_para(tf, [R(text, size, txt, bold=bold, font=font)], align=PP_ALIGN.CENTER, first=True)
    return sp

def status_badge(slide, x, y, text, kind='delivered', size=10.5):
    palette = {
        'delivered': (C.green, C.white),
        'ready':     (C.gold, C.ink),
        'pending':   (C.muted, C.white),
        'gov':       (C.navy, C.white),
        'live':      (C.success, C.white),
    }
    fill, txt = palette.get(kind, (C.teal, C.white))
    w = 0.16 + 0.085 * len(text)
    return chip(slide, x, y, w, 0.3, text, fill=fill, txt=txt, size=size, radius=0.5), w

def icon_badge(slide, cx, cy, d, fill=C.teal, glyphv='', glyph_color=C.white,
               glyph_size=16, ring=None):
    o = oval(slide, cx - d/2, cy - d/2, d, d, fill=fill)
    if ring:
        oval(slide, cx - d/2 - 0.04, cy - d/2 - 0.04, d + 0.08, d + 0.08, fill=None, line=ring, line_w=1.4)
    if glyphv:
        tf = o.text_frame; tf.word_wrap = False
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
        add_para(tf, [R(glyphv, glyph_size, glyph_color, bold=True)], align=PP_ALIGN.CENTER, first=True)
    return o

def najdi_star(slide, cx, cy, d, color=C.gold, alpha=None, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.STAR_8_POINT, Inches(cx - d/2), Inches(cy - d/2), Inches(d), Inches(d))
    no_shadow(sp)
    if alpha is not None:
        fill_alpha(sp, color, alpha)
    elif line:
        sp.fill.background(); sp.line.color.rgb = rgb(color); sp.line.width = Pt(line)
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = rgb(color)
    if not line:
        sp.line.fill.background()
    return sp

# ----------------------------------------------------------------- COMPOSITES
def process_flow(slide, x, y, w, h, steps, accent=C.teal, num_color=C.gold,
                 box_fill=C.white, connector_color=None, vertical=False):
    """Horizontal (default) or vertical chain of boxes joined by arrows.
    steps = list of (title, subtitle|None) or strings."""
    connector_color = connector_color or C.gold_dk
    n = len(steps)
    gap = (0.34 if not vertical else 0.26)
    if not vertical:
        bw = (w - gap * (n - 1)) / n
        bh = h
        for i, st in enumerate(steps):
            bx = x + i * (bw + gap)
            _flow_box(slide, bx, y, bw, bh, st, i + 1, accent, num_color, box_fill)
            if i < n - 1:
                ay = y + bh / 2
                connector(slide, bx + bw + 0.02, ay, bx + bw + gap - 0.02, ay,
                          color=connector_color, w=2.2, head=True)
    else:
        bh = (h - gap * (n - 1)) / n
        bw = w
        for i, st in enumerate(steps):
            by = y + i * (bh + gap)
            _flow_box(slide, x, by, bw, bh, st, i + 1, accent, num_color, box_fill)
            if i < n - 1:
                ax = x + bw / 2
                connector(slide, ax, by + bh + 0.02, ax, by + bh + gap - 0.02,
                          color=connector_color, w=2.2, head=True)
    return dict(x=x, y=y, w=w, h=h)

def _flow_box(slide, x, y, w, h, st, num, accent, num_color, box_fill):
    title, sub = (st if isinstance(st, tuple) else (st, None))
    card = rect(slide, x, y, w, h, fill=box_fill, line=C.line, line_w=1.0, radius=0.10, shadow=True)
    rect(slide, x, y, 0.07, h, fill=accent, radius=None)  # left accent
    # number bubble
    d = min(0.42, h * 0.42)
    icon_badge(slide, x + 0.30, y + 0.32, d, fill=accent, glyphv=str(num), glyph_color=num_color, glyph_size=15)
    tf, _ = textbox(slide, x + 0.16, y + 0.56, w - 0.3, h - 0.62, anchor=MSO_ANCHOR.TOP)
    add_para(tf, [R(title, 13.5, C.ink, bold=True, font=BODY)], line_spacing=1.0, first=True)
    if sub:
        add_para(tf, [R(sub, 10.5, C.muted)], space_before=2, line_spacing=1.0)
    return card

def chevron_flow(slide, x, y, w, h, steps, fills=None):
    """Overlapping chevrons (pentagon arrows) for a linear pipeline."""
    n = len(steps)
    overlap = 0.18
    sw_ = (w + overlap * (n - 1)) / n
    fills = fills or [C.teal_dk, C.teal, C.teal_md, C.green, C.gold_dk, C.gold]
    for i, st in enumerate(steps):
        sx = x + i * (sw_ - overlap)
        mso = MSO_SHAPE.PENTAGON if i == 0 else MSO_SHAPE.CHEVRON
        sp = freeshape(slide, mso, sx, y, sw_, h, fill=fills[i % len(fills)])
        tf = sp.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        tf.margin_left = Inches(0.18 if i == 0 else 0.26); tf.margin_right = Inches(0.1)
        txtcol = C.white if i < n - 1 or fills[i % len(fills)] != C.gold else C.ink
        add_para(tf, [R(st, 12.5, txtcol, bold=True)], align=PP_ALIGN.CENTER, first=True)
    return dict(x=x, y=y, w=w, h=h)

def hub_spoke(slide, cx, cy, hub_label, spokes, hub_d=1.7, spoke_d=1.55, radius=2.5,
              hub_fill=C.teal_dk, hub_sub=None, start_deg=-90):
    """Center hub with radial spoke cards. spokes=list of dict(title, sub, color, status)."""
    n = len(spokes)
    # connectors first (under nodes)
    pts = []
    for i in range(n):
        ang = math.radians(start_deg + i * 360.0 / n)
        sx = cx + radius * math.cos(ang)
        sy = cy + radius * math.sin(ang) * 0.78  # vertical squeeze
        pts.append((sx, sy, ang))
        connector(slide, cx, cy, sx, sy, color=C.line_dk, w=1.6, head=False)
    # hub
    grad = oval(slide, cx - hub_d/2, cy - hub_d/2, hub_d, hub_d)
    gradient_fill(grad, C.teal_dk, C.teal_md, 60); soft_shadow(grad, alpha=40)
    oval(slide, cx - hub_d/2 - 0.06, cy - hub_d/2 - 0.06, hub_d + 0.12, hub_d + 0.12, fill=None, line=C.gold, line_w=1.5)
    tf = grad.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
    add_para(tf, [R(hub_label, 15, C.white, bold=True, font=HEAD)], align=PP_ALIGN.CENTER, first=True)
    if hub_sub:
        add_para(tf, [R(hub_sub, 9.5, C.gold_lt)], align=PP_ALIGN.CENTER, space_before=1)
    # spokes
    for (sx, sy, ang), sp in zip(pts, spokes):
        col = sp.get('color', C.teal)
        node = oval(slide, sx - spoke_d/2, sy - spoke_d/2, spoke_d, spoke_d, fill=C.white, line=col, line_w=2.4)
        soft_shadow(node, alpha=26)
        tf = node.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
        tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
        add_para(tf, [R(sp['title'], 11.5, C.ink, bold=True)], align=PP_ALIGN.CENTER, first=True)
        if sp.get('sub'):
            add_para(tf, [R(sp['sub'], 9, C.muted)], align=PP_ALIGN.CENTER, space_before=1)
        if sp.get('status'):
            badge_kind = sp['status']
            txtmap = {'delivered': 'BUILT', 'ready': 'READY', 'pending': 'SANDBOX', 'gov': 'GOV-GATED'}
            chip(slide, sx - 0.42, sy + spoke_d/2 - 0.30, 0.84, 0.22,
                 txtmap.get(badge_kind, badge_kind), fill=col, txt=C.white, size=8, radius=0.5)
    return dict(x=cx-radius-spoke_d, y=cy-radius, w=2*(radius+spoke_d), h=2*radius)

def layered_stack(slide, x, y, w, h, layers, label_w=2.1):
    """Architecture layers stacked vertically. layers = list of dict(
       name, items=[...], fill, txt, tag=None). Top layer drawn first (at y)."""
    n = len(layers)
    gap = 0.12
    lh = (h - gap * (n - 1)) / n
    for i, L in enumerate(layers):
        ly = y + i * (lh + gap)
        fill = L.get('fill', C.card)
        txt = L.get('txt', C.ink)
        band = rect(slide, x, ly, w, lh, fill=fill, radius=0.06, shadow=True, line=L.get('line'), line_w=1.0)
        # label zone
        tf, _ = textbox(slide, x + 0.18, ly, label_w, lh, anchor=MSO_ANCHOR.MIDDLE)
        add_para(tf, [R(L['name'], 12.5, txt, bold=True, font=BODY)], line_spacing=1.0, first=True)
        if L.get('tag'):
            add_para(tf, [R(L['tag'], 8.5, txt)], space_before=1)
        # items as pills
        items = L.get('items', [])
        if items:
            ix = x + label_w + 0.25
            avail = w - label_w - 0.45
            pw = (avail - 0.12 * (len(items) - 1)) / len(items)
            pill_fill = L.get('pill', C.white)
            pill_txt = L.get('pill_txt', C.ink)
            for j, it in enumerate(items):
                px = ix + j * (pw + 0.12)
                ph = min(lh - 0.22, 0.46)
                pcard = rect(slide, px, ly + (lh - ph)/2, pw, ph, fill=pill_fill, radius=0.18, line=L.get('pill_line', C.line), line_w=0.75)
                ptf = pcard.text_frame; ptf.vertical_anchor = MSO_ANCHOR.MIDDLE; ptf.word_wrap = True
                ptf.margin_left = Inches(0.05); ptf.margin_right = Inches(0.05)
                add_para(ptf, [R(it, 9.8, pill_txt, bold=False)], align=PP_ALIGN.CENTER, first=True)
    return dict(x=x, y=y, w=w, h=h)

def comparison_table(slide, x, y, w, h, headers, rows, highlight_col=None,
                     col_widths=None, header_fill=C.teal_dk, first_col_fill=C.card2):
    """Styled comparison table with auto check/cross coloring.
    Cells '✓' green, '✗' muted-red, '~' gold. highlight_col tints a column."""
    ncols = len(headers); nrows = len(rows) + 1
    gf = slide.shapes.add_table(nrows, ncols, Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = gf.table
    # strip themed style for full control
    tblPr = tbl._tbl.find(qn('a:tblPr'))
    if tblPr is not None:
        for s in list(tblPr.findall(qn('a:tableStyleId'))):
            tblPr.remove(s)
        tblPr.set('firstRow', '0'); tblPr.set('bandRow', '0')
    # widths
    if col_widths:
        tot = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Inches(w * cw / tot)
    else:
        first = w * 0.34
        rest = (w - first) / (ncols - 1)
        tbl.columns[0].width = Inches(first)
        for i in range(1, ncols):
            tbl.columns[i].width = Inches(rest)
    hh = 0.52
    tbl.rows[0].height = Inches(hh)
    rh = (h - hh) / (nrows - 1)
    for r in range(1, nrows):
        tbl.rows[r].height = Inches(rh)

    def style_cell(cell, text, fill, txtcolor, bold=False, size=12, align=PP_ALIGN.LEFT, font=BODY):
        cell.fill.solid(); cell.fill.fore_color.rgb = rgb(fill)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
        tf = cell.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = align
        _apply_runs(p, [R(text, size, txtcolor, bold=bold, font=font)])

    # header row
    for c in range(ncols):
        hl = (highlight_col == c)
        style_cell(tbl.cell(0, c), headers[c], C.gold if hl else header_fill,
                   C.ink if hl else C.white, bold=True, size=12.5,
                   align=(PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER), font=BODY)
    # body
    for ri, row in enumerate(rows):
        r = ri + 1
        for c in range(ncols):
            val = str(row[c])
            hl = (highlight_col == c)
            base = C.gold_pale if hl else (C.white if r % 2 else C.card)
            if c == 0:
                style_cell(tbl.cell(r, c), val, first_col_fill, C.ink, bold=True, size=11.5, align=PP_ALIGN.LEFT)
                continue
            txtcol = C.ink_soft
            disp = val; bold = False; align = PP_ALIGN.CENTER; size = 11.5
            if val.strip() in ('✓', 'yes', 'Yes'):
                disp = '✓'; txtcol = C.success; bold = True; size = 15
            elif val.strip() in ('✗', 'no', 'No'):
                disp = '✗'; txtcol = C.crit; bold = True; size = 15
            elif val.strip() in ('~', 'partial', 'Partial'):
                disp = '◑'; txtcol = C.warn; bold = True; size = 14
            style_cell(tbl.cell(r, c), disp, base, txtcol, bold=bold, size=size, align=align)
    return dict(x=x, y=y, w=w, h=h, table=tbl)

def quadrant(slide, x, y, w, h, xlabel, ylabel, items, x_axis=('', ''), y_axis=('', '')):
    """2x2 positioning map. items=list of dict(label, qx(0-1), qy(0-1), color, winner)."""
    panel = rect(slide, x, y, w, h, fill=C.white, line=C.line, line_w=1.0, radius=0.04, shadow=True)
    # quadrant tints
    midx, midy = x + w/2, y + h/2
    fill_alpha(rect(slide, midx, y, w/2, h/2, radius=0.0), C.gold, 12)  # top-right highlight
    # axes
    connector(slide, x + 0.15, midy, x + w - 0.15, midy, color=C.line_dk, w=1.4, head=True, tail=True)
    connector(slide, midx, y + h - 0.15, midx, y + 0.15, color=C.line_dk, w=1.4, head=True, tail=True)
    # axis labels
    label(slide, x, y + h + 0.05, w, 0.3, xlabel, 10.5, C.muted, italic=True, align=PP_ALIGN.CENTER)
    ytf, ytb = textbox(slide, (x - 0.44) - h/2, (y + h/2) - 0.17, h, 0.34, anchor=MSO_ANCHOR.MIDDLE)
    ytf.word_wrap = False
    add_para(ytf, [R(ylabel, 10.5, C.muted, italic=True)], align=PP_ALIGN.CENTER, first=True)
    ytb.rotation = 270
    # corner hints
    if x_axis[0]: label(slide, x + 0.1, midy + 0.02, w/2 - 0.2, 0.25, x_axis[0], 8.5, C.muted_lt)
    if x_axis[1]: label(slide, midx + 0.1, midy + 0.02, w/2 - 0.2, 0.25, x_axis[1], 8.5, C.muted_lt, align=PP_ALIGN.RIGHT)
    # plot items
    for it in items:
        px = x + 0.4 + it['qx'] * (w - 0.8)
        py = y + h - 0.4 - it['qy'] * (h - 0.8)
        win = it.get('winner')
        d = 0.30 if win else 0.20
        col = it.get('color', C.muted)
        dot = oval(slide, px - d/2, py - d/2, d, d, fill=col, line=(C.gold if win else None), line_w=2)
        if win: soft_shadow(dot, alpha=40)
        lab = it['label']
        tf, _ = textbox(slide, px - 1.1, py + d/2 + 0.02, 2.2, 0.5, anchor=MSO_ANCHOR.TOP)
        add_para(tf, [R(lab, 10 if win else 9.5, C.ink if win else C.muted, bold=bool(win))],
                 align=PP_ALIGN.CENTER, first=True, line_spacing=0.95)
    return dict(x=x, y=y, w=w, h=h)

def timeline(slide, x, y, w, phases, line_y=None, accent=C.gold):
    """Horizontal milestone timeline. phases=list of dict(label, items=[...], color, tag)."""
    n = len(phases)
    line_y = line_y or (y + 0.0)
    connector(slide, x, line_y, x + w, line_y, color=C.line_dk, w=3, head=True)
    seg = w / n
    for i, ph in enumerate(phases):
        cx = x + seg * i + seg/2
        col = ph.get('color', C.teal)
        # node
        oval(slide, cx - 0.13, line_y - 0.13, 0.26, 0.26, fill=col, line=C.white, line_w=2)
        above = (i % 2 == 0)
        # tag chip
        if ph.get('tag'):
            chip(slide, cx - 0.55, line_y - 0.55, 1.1, 0.3, ph['tag'], fill=col, txt=C.white, size=9.5, radius=0.5)
        # card
        ch = 1.7
        cy = (line_y - 0.75 - ch) if above else (line_y + 0.55)
        card = rect(slide, cx - seg/2 + 0.12, cy, seg - 0.24, ch, fill=C.white, line=C.line, line_w=1, radius=0.08, shadow=True)
        rect(slide, cx - seg/2 + 0.12, cy, seg - 0.24, 0.08, fill=col)
        tf, _ = textbox(slide, cx - seg/2 + 0.28, cy + 0.16, seg - 0.5, ch - 0.28)
        add_para(tf, [R(ph['label'], 12.5, C.ink, bold=True)], first=True, line_spacing=1.0)
        for it in ph.get('items', []):
            add_para(tf, [R('•  ', 10, accent, bold=True), R(it, 10, C.muted)], space_before=4, line_spacing=0.98)
        # connector node->card
        connector(slide, cx, (cy + ch) if above else (line_y + 0.13),
                  cx, (line_y - 0.13) if above else cy, color=C.line_dk, w=1.2, head=False)
    return dict(x=x, y=y - 2.4, w=w, h=4.8)

def kpi_tiles(slide, x, y, w, h, tiles, cols=None, accent=C.gold,
              tile_fill=C.white, num_color=C.teal_dk):
    """Grid of big-number tiles. tiles=list of dict(num, label, sub=None, color=None)."""
    n = len(tiles)
    cols = cols or n
    rows = math.ceil(n / cols)
    gx, gy = 0.22, 0.22
    tw = (w - gx * (cols - 1)) / cols
    th = (h - gy * (rows - 1)) / rows
    for i, t in enumerate(tiles):
        r, c = divmod(i, cols)
        tx = x + c * (tw + gx)
        ty = y + r * (th + gy)
        card = rect(slide, tx, ty, tw, th, fill=tile_fill, line=C.line, line_w=1, radius=0.10, shadow=True)
        rect(slide, tx, ty, tw, 0.08, fill=t.get('color', accent))
        tf, _ = textbox(slide, tx + 0.12, ty + 0.10, tw - 0.24, th - 0.20, anchor=MSO_ANCHOR.MIDDLE)
        add_para(tf, [R(str(t['num']), min(34, th*36), t.get('num_color', num_color), bold=True, font=NUMS)],
                 align=PP_ALIGN.CENTER, first=True, line_spacing=0.9)
        add_para(tf, [R(t['label'], 11, C.ink_soft, bold=True)], align=PP_ALIGN.CENTER, space_before=2, line_spacing=0.95)
        if t.get('sub'):
            add_para(tf, [R(t['sub'], 9, C.muted)], align=PP_ALIGN.CENTER, space_before=1, line_spacing=0.95)
    return dict(x=x, y=y, w=w, h=h)

def capability_map(slide, x, y, w, h, groups, cols=4):
    """Capability grid grouped+color-coded. groups=list of dict(name,color,items=[...]).
    Renders a header pill per group with its tiles in a band."""
    n = len(groups)
    gap = 0.16
    gh = (h - gap * (n - 1)) / n
    for gi, g in enumerate(groups):
        gy = y + gi * (gh + gap)
        col = g['color']
        # group label
        chip(slide, x, gy, 1.95, gh, g['name'], fill=col, txt=C.white, size=11.5, radius=0.08)
        items = g['items']
        ix = x + 2.1
        avail = w - 2.1
        per = min(cols, len(items))
        # wrap into rows within the band
        rows = math.ceil(len(items) / per)
        tw = (avail - 0.12 * (per - 1)) / per
        ih = (gh - 0.10 * (rows - 1)) / rows
        for j, it in enumerate(items):
            r, c = divmod(j, per)
            tx = ix + c * (tw + 0.12)
            ty = gy + r * (ih + 0.10)
            tcard = rect(slide, tx, ty, tw, ih, fill=C.white, line=col, line_w=1.0, radius=0.10)
            tf = tcard.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
            tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
            add_para(tf, [R(it, 10, C.ink_soft, bold=True)], align=PP_ALIGN.CENTER, first=True, line_spacing=0.95)
    return dict(x=x, y=y, w=w, h=h)

def funnel(slide, x, y, w, h, stages):
    """Vertical narrowing funnel via freeform trapezoids (upright text)."""
    n = len(stages)
    gap = 0.12
    sh = (h - gap * (n - 1)) / n
    top_frac, min_frac = 1.0, 0.40
    for i, s in enumerate(stages):
        ft = top_frac - (top_frac - min_frac) * (i / n)
        fb = top_frac - (top_frac - min_frac) * ((i + 1) / n)
        wt, wb = w * ft, w * fb
        sy = y + i * (sh + gap)
        cx = x + w / 2
        pts = [(cx - wt/2, sy), (cx + wt/2, sy), (cx + wb/2, sy + sh), (cx - wb/2, sy + sh)]
        seg = poly(slide, pts, fill=s.get('color', C.teal))
        soft_shadow(seg, alpha=22)
        tf, _ = textbox(slide, cx - wt/2, sy, wt, sh, anchor=MSO_ANCHOR.MIDDLE)
        add_para(tf, [R(s['label'], 12.5, C.white, bold=True)], align=PP_ALIGN.CENTER, first=True)
        if s.get('sub'):
            add_para(tf, [R(s['sub'], 9.5, C.gold_lt)], align=PP_ALIGN.CENTER, space_before=1)
    return dict(x=x, y=y, w=w, h=h)

def ladder(slide, x, y, w, h, rungs, accent=C.gold):
    """Escalation ladder (ascending steps). rungs=list of (title, sub)."""
    n = len(rungs)
    gap = 0.16
    rh = (h - gap * (n - 1)) / n
    step_w = w * 0.6
    for i, rg in enumerate(reversed(rungs)):
        idx = n - 1 - i
        ry = y + i * (rh + gap)
        rw = step_w + (w - step_w) * (i / max(1, n - 1))
        card = rect(slide, x, ry, rw, rh, fill=(C.teal_dk if i == 0 else C.white),
                    line=accent if i == 0 else C.line, line_w=1.2, radius=0.10, shadow=True)
        icon_badge(slide, x + 0.34, ry + rh/2, min(0.5, rh*0.6), fill=accent, glyphv=str(idx + 1),
                   glyph_color=C.ink, glyph_size=14)
        title, sub = (rg if isinstance(rg, tuple) else (rg, None))
        tf, _ = textbox(slide, x + 0.7, ry, rw - 0.85, rh, anchor=MSO_ANCHOR.MIDDLE)
        add_para(tf, [R(title, 12.5, C.white if i == 0 else C.ink, bold=True)], first=True, line_spacing=1.0)
        if sub:
            add_para(tf, [R(sub, 9.5, C.gold_lt if i == 0 else C.muted)], space_before=1, line_spacing=0.95)
        if i < n - 1:
            connector(slide, x + rw*0.5, ry - 0.02, x + rw*0.5, ry - gap + 0.02, color=accent, w=2, head=True)
    return dict(x=x, y=y, w=w, h=h)

def gauge(slide, cx, cy, r, pct, label_txt='', color=C.teal, track=C.card3, ring_w=0.26):
    """Donut gauge showing pct (0-100)."""
    # track
    oval(slide, cx - r, cy - r, 2*r, 2*r, fill=track)
    # arc via pie shape
    arc = slide.shapes.add_shape(MSO_SHAPE.PIE, Inches(cx - r), Inches(cy - r), Inches(2*r), Inches(2*r))
    no_shadow(arc); arc.fill.solid(); arc.fill.fore_color.rgb = rgb(color); arc.line.fill.background()
    # pie adjustments: start/end angle (degrees). adj1=start, adj2=end
    try:
        arc.adjustments[0] = -90
        arc.adjustments[1] = -90 + 360 * (pct / 100.0)
    except Exception:
        pass
    # inner hole
    ir = r - ring_w
    oval(slide, cx - ir, cy - ir, 2*ir, 2*ir, fill=C.white)
    tf, _ = textbox(slide, cx - r, cy - r*0.55, 2*r, r, anchor=MSO_ANCHOR.MIDDLE)
    add_para(tf, [R('%d%%' % pct, r*36, C.teal_dk, bold=True, font=NUMS)], align=PP_ALIGN.CENTER, first=True, line_spacing=0.85)
    if label_txt:
        add_para(tf, [R(label_txt, 10, C.muted, bold=True)], align=PP_ALIGN.CENTER, space_before=1)
    return dict(x=cx-r, y=cy-r, w=2*r, h=2*r)

def progress_bar(slide, x, y, w, label_txt, pct, color=C.teal, h=0.30, show_pct=True):
    label(slide, x, y - 0.02, w, 0.26, label_txt, 11, C.ink, bold=True)
    by = y + 0.26
    rect(slide, x, by, w, h, fill=C.card3, radius=0.5)
    rect(slide, x, by, max(0.12, w * pct / 100.0), h, fill=color, radius=0.5)
    if show_pct:
        label(slide, x + w - 0.7, by - 0.02, 0.7, h + 0.04, '%d%%' % pct, 10.5, C.teal_dk, bold=True,
              align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    return dict(x=x, y=y, w=w, h=h + 0.26)

# ----------------------------------------------------------------- images
def _png_size(path):
    try:
        with open(path, 'rb') as f:
            head = f.read(26)
        if head[:8] == b'\x89PNG\r\n\x1a\n':
            w, h = struct.unpack('>II', head[16:24])
            return w, h
    except Exception:
        pass
    return None

def picture_cover(slide, x, y, w, h, path, bias_top=True):
    """Place image filling x,y,w,h with center/top-biased crop (no distortion)."""
    dims = _png_size(path)
    pic = slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    if dims:
        iw, ih = dims
        target = w / h
        ia = iw / ih
        if ia > target:  # too wide -> crop sides
            crop = 1 - target / ia
            pic.crop_left = crop / 2; pic.crop_right = crop / 2
        else:            # too tall -> crop vertical
            crop = 1 - ia / target
            if bias_top:
                pic.crop_top = 0.0; pic.crop_bottom = crop
            else:
                pic.crop_top = crop / 2; pic.crop_bottom = crop / 2
    return pic

def browser_frame(slide, x, y, w, h, img_path, title='', accent=C.gold, dark=True):
    """Browser/device chrome around a screenshot. Returns bbox."""
    bezel = rect(slide, x, y, w, h, fill=(C.teal_dk2 if dark else C.card), radius=0.05, shadow=True)
    soft_shadow(bezel, blur=0.16, dist=0.08, alpha=40)
    bar_h = 0.34
    rect(slide, x, y, w, bar_h, fill=(C.teal_dk if dark else C.card2), radius=0.05)
    rect(slide, x, y + bar_h - 0.06, w, 0.06, fill=(C.teal_dk if dark else C.card2))
    for i, cc in enumerate(['C24233', 'D79A2B', '1F9D6B']):
        oval(slide, x + 0.16 + i * 0.16, y + bar_h/2 - 0.045, 0.09, 0.09, fill=cc)
    # url pill
    pill = rect(slide, x + 0.66, y + 0.06, min(w - 1.0, 4.2), bar_h - 0.12, fill=(C.teal_dk2 if dark else C.white), radius=0.5)
    ptf = pill.text_frame; ptf.vertical_anchor = MSO_ANCHOR.MIDDLE; ptf.word_wrap = False
    ptf.margin_left = Inches(0.12)
    add_para(ptf, [R('🔒  ' + (title or 'watheeq.clario360.sa'), 9, C.gold_lt if dark else C.muted)], first=True)
    pad = 0.06
    if img_path and os.path.exists(img_path):
        picture_cover(slide, x + pad, y + bar_h, w - 2*pad, h - bar_h - pad, img_path, bias_top=True)
    else:
        rect(slide, x + pad, y + bar_h, w - 2*pad, h - bar_h - pad, fill=C.card)
    return dict(x=x, y=y, w=w, h=h)

# ----------------------------------------------------------------- scaffolds
def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _footer(slide, num):
    label(slide, 0.62, SH - 0.46, 7.0, 0.3, 'Watheeq (وثيق)   ·   Clario360 Sovereign Platform', 9, C.muted_lt)
    label(slide, SW - 1.7, SH - 0.46, 1.1, 0.3, '%02d / 40' % num, 9.5, C.muted, bold=True, align=PP_ALIGN.RIGHT)
    try:
        slide.shapes.add_picture(LOGO_ICON, Inches(SW - 0.52), Inches(SH - 0.52), Inches(0.30), Inches(0.30))
    except Exception:
        pass

def content_scaffold(prs, num, title, subtitle=None, kicker=None):
    """Standard content slide chrome. Returns (slide, content_rect dict)."""
    slide = new_slide(prs)
    rect(slide, -0.1, -0.1, SW + 0.2, SH + 0.2, fill=C.bg)
    rect(slide, 0, 0, 0.16, SH, fill=C.gold)              # left accent edge
    najdi_star(slide, SW - 0.62, 0.62, 0.5, color=C.gold, alpha=12)  # faint corner motif
    # number chip
    nb = rect(slide, 0.62, 0.46, 0.92, 0.66, fill=C.teal_dk, radius=0.16, shadow=True)
    ntf = nb.text_frame; ntf.vertical_anchor = MSO_ANCHOR.MIDDLE; ntf.word_wrap = False
    add_para(ntf, [R('%02d' % num, 23, C.gold, bold=True, font=NUMS)], align=PP_ALIGN.CENTER, first=True)
    # kicker + title
    tx = 1.78
    if kicker:
        label(slide, tx, 0.40, 10.6, 0.26, kicker.upper(), 10.5, C.gold_dk, bold=True, spacing=1.5)
        ty = 0.62
    else:
        ty = 0.50
    label(slide, tx, ty, 10.9, 0.8, title, 26, C.teal_dk, bold=True, font=HEAD, line_spacing=0.98)
    if subtitle:
        label(slide, tx, ty + 0.62, 10.9, 0.4, subtitle, 14, C.muted, italic=True, font=BODY)
    rect(slide, tx, 1.66, 10.9, 0.032, fill=C.gold)       # gold rule
    _footer(slide, num)
    crect = dict(x=0.62, y=1.92, w=12.1, h=4.55)
    return slide, crect

def cover_slide(prs, title, subtitle=None, bullets_list=None, kicker=None):
    slide = new_slide(prs)
    gradient_fill(rect(slide, -0.1, -0.1, SW + 0.2, SH + 0.2), C.teal_dk2, C.teal_dk, 120)
    rect(slide, 0, 0, SW, 0.16, fill=C.gold)
    rect(slide, 0, SH - 0.16, SW, 0.16, fill=C.gold)
    # motif cluster
    najdi_star(slide, SW - 1.9, SH - 1.7, 2.7, color=C.gold, alpha=8)
    najdi_star(slide, SW - 1.0, 1.2, 1.1, color=C.gold, alpha=10)
    for i in range(3):
        najdi_star(slide, 11.0 + i*0.42, 0.78, 0.26, color=C.gold, line=1.0)
    # logo
    try:
        slide.shapes.add_picture(LOGO_ICON, Inches(0.95), Inches(0.85), Inches(0.95), Inches(0.95))
    except Exception:
        pass
    if kicker:
        label(slide, 2.05, 1.06, 9.5, 0.4, kicker.upper(), 13, C.gold, bold=True, spacing=2)
    label(slide, 0.95, 2.35, 11.4, 2.3, title, 42, C.white, bold=True, font=HEAD, line_spacing=1.02)
    if subtitle:
        label(slide, 0.97, 4.78, 11.0, 0.9, subtitle, 20, C.gold_lt, italic=True)
    if bullets_list:
        tf, _ = textbox(slide, 0.97, 5.55, 11.4, 1.3)
        for i, b in enumerate(bullets_list[:3]):
            add_para(tf, [R('◆  ', 12, C.gold), R(b, 14.5, 'DCE8E5')], space_after=5, first=(i == 0))
    label(slide, 0.97, SH - 0.66, 11.4, 0.34, 'Watheeq (وثيق) — trustworthy · reliable', 12, C.gold)
    return slide

def section_slide(prs, num, title, subtitle=None):
    slide = new_slide(prs)
    gradient_fill(rect(slide, -0.1, -0.1, SW + 0.2, SH + 0.2), C.teal_dk, C.teal, 120)
    rect(slide, 0, 0, 0.16, SH, fill=C.gold)
    najdi_star(slide, SW - 1.6, SH/2, 2.4, color=C.gold, alpha=9)
    label(slide, 1.0, SH/2 - 1.0, 1.6, 0.8, '%02d' % num, 60, C.gold, bold=True, font=NUMS)
    label(slide, 1.0, SH/2 - 0.05, 11.0, 1.2, title, 34, C.white, bold=True, font=HEAD)
    if subtitle:
        label(slide, 1.0, SH/2 + 0.85, 11.0, 0.6, subtitle, 16, C.gold_lt, italic=True)
    return slide

def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text or ''

# ----------------------------------------------------------------- md parser
def parse_md(path):
    text = open(path, encoding='utf-8').read()
    chunks = re.split(r"(?m)^## Slide\s+(\d+)\s+—\s+(.+)$", text)
    slides = []
    for i in range(1, len(chunks), 3):
        num = int(chunks[i]); title = chunks[i+1].strip(); body = chunks[i+2].split('\n---')[0]
        paras = [p.strip() for p in re.split(r"\n\s*\n", body) if p.strip()]
        subtitle, items, visual, notes = '', [], '', ''
        for p in paras:
            first = p.splitlines()[0].strip()
            if first.startswith('- '):
                for ln in p.splitlines():
                    ln = ln.strip()
                    if ln.startswith('- '):
                        items.append(_strip(ln[2:]))
            elif first.startswith('**Visual:**'):
                visual = _strip(' '.join(l.strip() for l in p.splitlines()).replace('**Visual:**', ''))
            elif first.startswith('**Speaker notes:**'):
                notes = _strip(' '.join(l.strip() for l in p.splitlines()).replace('**Speaker notes:**', ''))
            elif re.match(r'^\*\*.+\*\*$', first) and not subtitle:
                subtitle = _strip(first)
        slides.append(dict(num=num, title=title, subtitle=subtitle, bullets=items, visual=visual, notes=notes))
    return slides

def _strip(s):
    s = re.sub(r'\*\*(.+?)\*\*', r'\1', s)
    s = re.sub(r'\*(.+?)\*', r'\1', s)
    return s.strip()

# ----------------------------------------------------------------- qa
def audit(prs, eps_in=0.18):
    eps = eps_in * EMU_IN
    issues = []
    for i, slide in enumerate(prs.slides, 1):
        for s in slide.shapes:
            try:
                l, t, w, h = s.left, s.top, s.width, s.height
                if l is None:
                    continue
            except Exception:
                continue
            if l < -eps or t < -eps or (l + w) > prs.slide_width + eps or (t + h) > prs.slide_height + eps:
                issues.append((i, str(s.shape_type), round(l/EMU_IN, 2), round(t/EMU_IN, 2),
                               round(w/EMU_IN, 2), round(h/EMU_IN, 2)))
    return issues

def render_pngs(pptx_path, outdir, dpi=150):
    soffice = '/Applications/LibreOffice.app/Contents/MacOS/soffice'
    if not os.path.exists(soffice):
        return None
    os.makedirs(outdir, exist_ok=True)
    pdf = os.path.join(outdir, os.path.splitext(os.path.basename(pptx_path))[0] + '.pdf')
    subprocess.run([soffice, '--headless', '--convert-to', 'pdf', '--outdir', outdir, pptx_path],
                   capture_output=True, timeout=300)
    if not os.path.exists(pdf):
        return None
    # pdf -> png via sips per page? sips can't split pdf pages. use pdftoppm if present, else 'sips' on pdf gives 1 page.
    pdftoppm = subprocess.run(['which', 'pdftoppm'], capture_output=True, text=True).stdout.strip()
    if pdftoppm:
        subprocess.run([pdftoppm, '-png', '-r', str(dpi), pdf, os.path.join(outdir, 'slide')],
                       capture_output=True, timeout=300)
    return outdir
